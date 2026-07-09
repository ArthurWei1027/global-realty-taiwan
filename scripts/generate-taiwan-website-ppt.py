"""Generate July sales meeting PPT: Taiwan website sitemap & progress."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "presentations" / "Global-Realty-台灣官網進度-202607.pptx"

TEAL = RGBColor(0x2D, 0x7A, 0x7A)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
GOLD = RGBColor(0xB4, 0x53, 0x09)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)


def add_bullets(slide, items: list[str], left=0.7, top=1.35, width=8.8, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)
    return box


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, title)
    cols, row_count = len(headers), len(rows) + 1
    table = slide.shapes.add_table(row_count, cols, Inches(0.5), Inches(1.4), Inches(9.0), Inches(0.45 * row_count)).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = NAVY
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT


def add_timeline_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "專案時間線", "Project timeline · 預計 7 月底完成上線")

    milestones = [
        ("6/19", "開始製作", "專案啟動、需求與架構規劃", True),
        ("6/23", "初步 UI 完成", "7 頁靜態 UI、Sitemap、設計預覽站", True),
        ("6/27", "設計討論", "雙品牌策略、UI 方向對齊", True),
        ("7/2", "最終 UI 討論", "子分頁、Footer、設計稿定稿", True),
        ("7/3", "月會進度報告", "Global Realty 全体销售大会", True),
        ("7月底", "預計完成上線", "正式素材、表單、網域部署", False),
    ]

    y_line = 2.15
    x_start, x_end = 0.9, 9.1
    line = slide.shapes.add_shape(1, Inches(x_start), Inches(y_line), Inches(x_end - x_start), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    line.line.fill.background()

    count = len(milestones)
    for i, (date, title, desc, done) in enumerate(milestones):
        x = x_start + (x_end - x_start) * i / (count - 1)
        dot_size = 0.22
        dot = slide.shapes.add_shape(
            9, Inches(x - dot_size / 2), Inches(y_line - dot_size / 2 + 0.01), Inches(dot_size), Inches(dot_size)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL if done else GOLD
        dot.line.fill.background()

        date_box = slide.shapes.add_textbox(Inches(x - 0.55), Inches(y_line - 0.75), Inches(1.1), Inches(0.35))
        dp = date_box.text_frame.paragraphs[0]
        dp.text = date
        dp.font.size = Pt(13)
        dp.font.bold = True
        dp.font.color.rgb = TEAL if done else GOLD
        dp.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(x - 0.75), Inches(y_line + 0.35), Inches(1.5), Inches(0.45))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(11)
        tp.font.bold = True
        tp.font.color.rgb = NAVY
        tp.alignment = PP_ALIGN.CENTER

        desc_box = slide.shapes.add_textbox(Inches(x - 0.85), Inches(y_line + 0.78), Inches(1.7), Inches(1.0))
        df = desc_box.text_frame
        df.word_wrap = True
        dp2 = df.paragraphs[0]
        dp2.text = desc
        dp2.font.size = Pt(9)
        dp2.font.color.rgb = GRAY
        dp2.alignment = PP_ALIGN.CENTER

    legend = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(8.6), Inches(0.6))
    lp = legend.text_frame.paragraphs[0]
    lp.text = "● 已完成　　● 預計目標　　從 6/19 啟動至 7 月底，約 6 週完成台灣官網建置與上線"
    lp.font.size = Pt(12)
    lp.font.color.rgb = NAVY


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s1, NAVY)
    tbox = s1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(3.5))
    tf = tbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = [
        ("環球置業台灣官網建置進度", 36, True, WHITE),
        ("Global Realty Taiwan Website — Sitemap & Progress", 20, False, RGBColor(0x94, 0xA3, 0xB8)),
        ("", 12, False, WHITE),
        ("Global Realty 7月全体销售大会 · Part 1", 16, False, TEAL),
        ("Arthur · 2026年7月3日", 16, False, RGBColor(0xCB, 0xD5, 0xE1)),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    # Slide 2 — 一句話定位
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s2, WHITE)
    add_title_bar(s2, "專案定位", "Project positioning")
    add_bullets(
        s2,
        [
            "面向台灣置產受眾的繁體中文官網，以環球置業 Global Realty 為主體驗。",
            "澳華國際集團 Award Global 以「from」署名與子分頁呈現（類 Facebook from Meta）。",
            "目標：3 秒辨識環球置業 · 10 秒理解集團背書 · 置產轉換不失焦。",
            "台北 101 45 樓亞太運營中心為服務據點，串聯澳洲本地資源。",
        ],
        size=20,
    )

    # Slide 3 — Sitemap 視覺
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s3, WHITE)
    add_title_bar(s3, "網站架構 Sitemap", "Information architecture")
    sitemap_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.2))
    stf = sitemap_box.text_frame
    tree = """頂層導覽（6 項）
  01 首頁 index.html
  02 精選建案 properties.html
  03 租賃管理 leasing.html
  04 活動預告 events.html
  05 澳洲不動產小課堂 classroom.html
  06 關於我們 about.html
        ├─ 關於環球置業 about.html（子分頁）
        └─ 澳華國際集團 group.html（子分頁）

其他：sitemap.html · sitemap.xml · robots.txt"""
    p = stf.paragraphs[0]
    p.text = tree
    p.font.name = "Consolas"
    p.font.size = Pt(16)
    p.font.color.rgb = NAVY

    note = s3.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.8))
    np = note.text_frame.paragraphs[0]
    np.text = "集團頁不佔頂層第 7 項，從 Header 下拉、頁頂 tab、from 署名進入"
    np.font.size = Pt(13)
    np.font.color.rgb = GOLD
    np.font.italic = True

    # Slide 4 — 頁面表
    add_table_slide(
        prs,
        "頁面一覽與狀態",
        ["#", "頁面", "檔案", "狀態"],
        [
            ["01", "首頁", "index.html", "✅ UI 完成"],
            ["02", "精選建案", "properties.html", "✅ UI 完成"],
            ["03", "租賃管理", "leasing.html", "✅ UI 完成"],
            ["04", "活動預告", "events.html", "✅ UI 完成"],
            ["05", "小課堂", "classroom.html", "✅ UI 完成"],
            ["06", "關於環球置業", "about.html", "✅ UI 完成"],
            ["07", "澳華國際集團", "group.html", "✅ UI 完成"],
        ],
    )

    # Slide 5 — 功能亮點
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s5, WHITE)
    add_title_bar(s5, "已完成功能", "Key features delivered")
    add_bullets(
        s5,
        [
            "雙品牌 Header / Footer：環球置業主 Logo + from 澳華國際集團",
            "首頁品牌關係區：GR 數據、架構樹、了解集團 CTA",
            "建案頁：雪梨 / 墨爾本篩選、623 Collins / AURA 等代理建案",
            "諮詢表單、右側浮動聯絡 icon（電話 / Email / LINE / FB）",
            "Footer 版權列：Copyright 2026 · Global Realty · 隱私權政策",
            "7 頁 Desktop 設計稿 + Live 預覽站（design/）",
        ],
        size=17,
    )

    # Slide 6 — 設計稿預覽
    s6img = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s6img, WHITE)
    add_title_bar(s6img, "設計稿預覽（首頁）", "Desktop mockup — 1440px")
    img_path = Path(__file__).resolve().parent.parent / "design" / "desktop" / "01-index.png"
    if img_path.exists():
        s6img.shapes.add_picture(str(img_path), Inches(0.6), Inches(1.25), width=Inches(8.8))

    # Slide 7 — 品牌 UI
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s6, WHITE)
    add_title_bar(s6, "UI / 品牌設計重點", "Design highlights")
    add_bullets(
        s6,
        [
            "主色：青綠 #2d7a7a（環球置業）+ 深藍（專業感）",
            "Header：6 項導覽 +「關於我們」下拉子分頁",
            "about / group 頁頂子導覽 tab，切換不迷路",
            "group.html：六大業務板塊、集團數據、Tony Lam、GR 高亮",
            "文案主語統一「環球置業」；集團數據僅在集團頁",
            "Design Mode：?design=1 可現場 demo 給主管審閱",
        ],
        size=17,
    )

    # Slide 8 — 時間線
    add_timeline_slide(prs)

    # Slide 9 — 進度
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s7, WHITE)
    add_title_bar(s7, "整體進度", "Overall progress")
    add_bullets(
        s7,
        [
            "【已完成 ~85%】靜態 UI、Sitemap、7 頁設計稿、雙品牌架構",
            "【目標】2026 年 7 月底完成網站上線",
            "【進行中】正式圖片替換（hero / 建案渲染圖）、文案終審",
            "【待確認】正式網域 / canonical 策略（目前指向 awardglobal.com.au）",
            "【待接入】表單後端 CRM、GTM 追蹤、正式房源 API 同步",
            "【待提供】官方 Logo 向量檔（目前為結構化 SVG 占位）",
        ],
        size=17,
    )

    # Slide 10 — 下一步
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s8, WHITE)
    add_title_bar(s8, "下一步 & 需要支援", "Next steps")
    add_bullets(
        s8,
        [
            "7月：主管 / Josh 審閱設計稿，確認 UI 與雙品牌策略",
            "銷售團隊：提供 623 / AURA 最新素材與台灣活動內容",
            "IT：確認上線 host、表單收件、網域 DNS",
            "上線前：合規免責聲明終審、Privacy Policy 連結確認",
        ],
        size=19,
    )

    # Slide 11 — Demo
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s9, NAVY)
    tbox9 = s9.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(3))
    tf9 = tbox9.text_frame
    for i, (text, size, color) in enumerate(
        [
            ("Live 預覽（本地）", 32, TEAL),
            ("http://localhost:8765/design/", 22, WHITE),
            ("Design Mode：/index.html?design=1", 18, RGBColor(0xCB, 0xD5, 0xE1)),
            ("", 8, WHITE),
            ("Q & A  ·  謝謝", 28, GOLD),
        ]
    ):
        p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
